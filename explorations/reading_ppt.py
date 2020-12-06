from pptx import Presentation  
from pptx.enum.shapes import MSO_SHAPE_TYPE   

prs = Presentation("./test_powerpoint.pptx")

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = []


for slide in prs.slides:
    for shape in slide.shapes:
        try:
            print(shape.image._filename)
        except AttributeError:
            continue
        # if shape.shape_id == 6:
        #     shape.text = "I have been inserted with python"
        #     print(shape.text)

        # if not shape.has_text_frame:
        #     continue
        # for paragraph in shape.text_frame.paragraphs:
        #     for run in paragraph.runs:
        #         text_runs.append(run.text)

  
# slide = prs.slides.add_slide(prs.slide_layouts[11])
# prs.save("./test_powerpoint.pptx") 

print(text_runs)