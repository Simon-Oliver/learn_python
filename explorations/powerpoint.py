from pptx import Presentation

pp = Presentation("./test_powerpoint.pptx")

for slide in pp.slides:  
    for shape in slide.shapes:  
        if not shape.has_text_frame:  
            continue 
        for paragraph in shape.text_frame.paragraphs:  
            for run in paragraph.runs:  
                print(run.text) 


# pp.save("Output.pptx") 